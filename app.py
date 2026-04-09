'''
# ========== 最顶部：强制国内镜像配置 ==========
import os
os.environ['HF_ENDPOINT'] = 'https://hf-mirror.com'
os.environ['HF_HOME'] = './hf_cache'

import streamlit as st
import shutil
import re
import uuid
from datetime import datetime
from rag_agent import RAGAgent
from config import MODEL_NAME, TOP_K, DATA_DIR
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from pptx import Presentation
from vector_store import VectorStore

# ========== 导入拆分的模块 ==========
from exam_core import clean_option_prefix, grade_exam_answers, format_answer
from exam_ui import render_view_mode, render_test_mode, render_graded_results

# ========== 知识点分析模块 ==========
from knowledge_analysis import render_knowledge_analysis_page, extract_knowledge_points

# ========== Sentence-BERT 核心评分模块 ==========
from sentence_transformers import SentenceTransformer

# ========== 全局SSL配置 ==========
import ssl
ssl._create_default_https_context = ssl._create_unverified_context

# ========== 模型加载 ==========
@st.cache_resource
def _load_sbert_model_internal():
    try:
        return SentenceTransformer('sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2')
    except Exception as e:
        return None

if 'model_initialized' not in st.session_state:
    st.session_state.model_initialized = False
    st.session_state.model = None
if not st.session_state.model_initialized:
    with st.spinner("正在加载语义匹配模型（仅首次加载需要1-2分钟）..."):
        loaded_model = _load_sbert_model_internal()
        if loaded_model is None:
            st.error("模型加载失败，已自动降级为关键词匹配模式")
        st.session_state.model = loaded_model
        st.session_state.model_initialized = True
model = st.session_state.model

# ========== 工具函数 ==========
def save_exam_history(questions: list, q_type: str, knowledge_point: str, difficulty: str, agent):
    if "exam_history" not in st.session_state:
        st.session_state.exam_history = []
    
    if knowledge_point.strip():
        doc_title = knowledge_point
    else:
        try:
            doc_names = agent.vector_store.get_all_filenames()
            doc_title = "、".join(doc_names) if doc_names else "全部知识点"
        except:
            doc_title = "全部知识点"
    
    exam_id = str(uuid.uuid4())[:8]
    history_item = {
        "id": exam_id,
        "time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "doc_title": doc_title,
        "knowledge_point": knowledge_point,
        "q_type": q_type,
        "difficulty": difficulty,
        "questions": questions,
        "main_test_record": None,
        "practice_records": []
    }
    st.session_state.exam_history.insert(0, history_item)
    return exam_id

def load_ppt_native(file_path: str):
    prs = Presentation(file_path)
    docs = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        text_content = "\n".join(slide_text)
        if text_content.strip():
            docs.append({
                "page_content": text_content,
                "metadata": {"page": slide_num, "filename": os.path.basename(file_path)}
            })
    return docs

# ========== 页面配置 ==========
st.set_page_config(
    page_title="智能题库系统",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ========== 【舒适版】全局布局CSS ==========
st.markdown("""
<style>
/* ========== 1. 彻底隐藏侧边栏折叠按钮，消除顶部空隙 ========== */
[data-testid="stSidebarCollapseButton"] {
    display: none !important;
}
[data-testid="stSidebarHeader"] {
    display: none !important;
    height: 0 !important;
    padding: 0 !important;
    margin: 0 !important;
}

/* ========== 2. 侧边栏容器：舒适的边距 ========== */
[data-testid="stSidebar"] {
    padding: 0 !important;
    margin: 0 !important;
}
[data-testid="stSidebar"] > div:first-child {
    padding-top: 0.8rem !important; /* 顶部留一点舒适的边距 */
    padding-bottom: 0.8rem !important;
    padding-left: 1rem !important;
    padding-right: 1rem !important;
    margin: 0 !important;
    width: 100% !important;
    top: 0 !important;
    position: absolute !important;
    height: 100vh !important;
}

/* ========== 3. 侧边栏所有元素间距：舒适不拥挤 ========== */
[data-testid="stSidebar"] [data-testid="stVerticalBlock"] > div {
    gap: 0.4rem !important; /* 元素之间的间距调大，更舒适 */
    padding: 0 !important;
    margin: 0 !important;
}

/* ========== 4. 侧边栏标题通用样式 ========== */
.sidebar-title-main {
    margin: 0 0 0.2rem 0 !important; /* 标题下面留一点间距 */
    padding: 0 !important;
    line-height: 1.3 !important;
}
.sidebar-title-sub {
    margin: 0 0 0.15rem 0 !important;
    padding: 0 !important;
    line-height: 1.3 !important;
}
/* 【核心】功能导航标题专属样式：只在下面加更大的空隙 */
.sidebar-title-nav {
    margin: 0 0 0.6rem 0 !important; /* 功能导航和智能出题之间的空隙更大 */
    padding: 0 !important;
    line-height: 1.3 !important;
}

/* ========== 5. 侧边栏所有标题通用margin控制 ========== */
[data-testid="stSidebar"] h1, 
[data-testid="stSidebar"] h2, 
[data-testid="stSidebar"] h3,
[data-testid="stSidebar"] h4,
[data-testid="stSidebar"] h5,
[data-testid="stSidebar"] h6 {
    margin-top: 0 !important;
    margin-bottom: 0.15rem !important;
    padding: 0 !important;
    line-height: 1.3 !important;
}

/* ========== 6. 按钮：舒适的高度和边距 ========== */
[data-testid="stSidebar"] button {
    margin-top: 0.1rem !important;
    margin-bottom: 0.1rem !important;
    min-height: 2.2rem !important; /* 按钮高度稍微调大，更易点击 */
    padding-top: 0.35rem !important;
    padding-bottom: 0.35rem !important;
    line-height: 1.2 !important;
}

/* ========== 7. 文件上传组件：舒适的边距 + 中文替换 ========== */
[data-testid="stSidebar"] [data-testid="stFileUploader"] {
    margin-top: 0.15rem !important;
    margin-bottom: 0.15rem !important;
    padding: 0 !important;
}
[data-testid="stSidebar"] [data-testid="stFileUploader"] > section {
    padding: 0.6rem 0.5rem !important;
    min-height: auto !important;
    margin: 0 !important;
}

/* ========== 8. 已上传文件列表：舒适的边距 ========== */
[data-testid="stSidebar"] [data-testid="stFileUploaderFile"] {
    margin-top: 0.1rem !important;
    margin-bottom: 0.1rem !important;
    padding-top: 0.2rem !important;
    padding-bottom: 0.2rem !important;
    min-height: auto !important;
}

/* ========== 9. 小字说明：舒适的边距 ========== */
[data-testid="stSidebar"] small {
    margin-top: 0.05rem !important;
    margin-bottom: 0.05rem !important;
    line-height: 1.2 !important;
}

/* ========== 10. 分割线：舒适的边距 ========== */
[data-testid="stSidebar"] hr {
    margin-top: 0.3rem !important;
    margin-bottom: 0.3rem !important;
    padding: 0 !important;
}

/* ========== 11. 文件上传框替换为中文 ========== */
/* 替换主标题：Drag and drop files here */
[data-testid="stFileUploader"] section > div:first-child > span:first-child {
    font-size: 0 !important;
}
[data-testid="stFileUploader"] section > div:first-child > span:first-child::before {
    content: "拖拽文件到此处";
    font-size: 2.4rem;
    color: #262730;
}

/* 替换限制说明：Limit 200MB per file... */
[data-testid="stFileUploader"] section > div:first-child > small {
    font-size: 0 !important;
}
[data-testid="stFileUploader"] section > div:first-child > small::before {
    content: "单文件最大200MB • 支持 PDF, DOCX, TXT, PPT 等格式";
    font-size: 1.1rem;
    color: #6b7280;
}

/* 替换按钮文本：Browse files */
[data-testid="stFileUploader"] section button p {
    font-size: 0 !important;
}
[data-testid="stFileUploader"] section button p::before {
    content: "选择文件";
    font-size: 1.4rem;
}

/* 适配深色模式的颜色 */
@media (prefers-color-scheme: dark) {
    [data-testid="stFileUploader"] section > div:first-child > span:first-child::before {
        color: #fafafa;
    }
    [data-testid="stFileUploader"] section > div:first-child > small::before {
        color: #9ca3af;
    }
}

/* ========== 主页面正确优化 ========== */
[data-testid="stMainBlockContainer"] {
    padding-top: 2rem !important;
    padding-bottom: 2rem !important;
    max-width: 90% !important;
}
[data-testid="stMainBlockContainer"] h1 {
    margin-top: 0.3rem !important;
    margin-bottom: 0.8rem !important;
    line-height: 1.3 !important;
}
[data-testid="stMainBlockContainer"] h2, 
[data-testid="stMainBlockContainer"] h3 {
    margin-top: 0.7rem !important;
    margin-bottom: 0.4rem !important;
    line-height: 1.3 !important;
}
[data-testid="stMainBlockContainer"] hr {
    margin-top: 0.7rem !important;
    margin-bottom: 0.7rem !important;
}
[data-testid="stMainBlockContainer"] [data-testid="stTextInput"],
[data-testid="stMainBlockContainer"] [data-testid="stSelectbox"],
[data-testid="stMainBlockContainer"] [data-testid="stSlider"] {
    margin-top: 0.05rem !important;
    margin-bottom: 0.05rem !important;
}
[data-testid="stMainBlockContainer"] button {
    margin-top: 0.1rem !important;
    margin-bottom: 0.1rem !important;
}
</style>
""", unsafe_allow_html=True)

# ========== session_state 全量初始化 ==========
if "initialized" not in st.session_state:
    st.cache_data.clear()
    st.cache_resource.clear()
    st.session_state.initialized = True
    st.session_state.page_state = "main"
    st.session_state.exam_history = []
    st.session_state.generated_questions = []
    st.session_state.current_q_type = ""
    st.session_state.exam_mode = "view"
    st.session_state.user_answers = {}
    st.session_state.graded = False
    st.session_state.score_detail = []
    st.session_state.total_score = 0
    if "user_answer_records" not in st.session_state:
        st.session_state.user_answer_records = []
    st.session_state.history_current_exam_id = ""
    st.session_state.history_questions = []
    st.session_state.history_q_type = ""
    if "messages" not in st.session_state:
        st.session_state.messages = []
    st.session_state.wrong_questions = []
    st.session_state.wrong_exam_mode = "view"
    st.session_state.wrong_user_answers = {}
    st.session_state.wrong_graded = False
    st.session_state.wrong_total_score = 0
    st.session_state.wrong_score_detail = []
    st.session_state.selected_wrong_ids = []
    if "uploader_key" not in st.session_state:
        st.session_state.uploader_key = 0
    # 【新增】控制主页面是否显示导入成功提示
    if "show_import_success" not in st.session_state:
        st.session_state.show_import_success = False

# RAG Agent 初始化
if "rag_agent" not in st.session_state:
    st.session_state.rag_agent = RAGAgent(model=MODEL_NAME, top_k=TOP_K)

# ========== 【舒适版】左侧侧边栏 ==========
with st.sidebar:
    # 1. 系统名称
    st.markdown('<h4 class="sidebar-title-main">🧠 智能题库系统</h4>', unsafe_allow_html=True)
    st.divider()

    # 2. 文档导入区域
    st.markdown('<h5 class="sidebar-title-sub">📂 文档导入</h5>', unsafe_allow_html=True)
    st.caption("支持 PDF/DOCX/TXT/PPT/PPTX")
    # 文件上传组件
    uploaded_files = st.file_uploader(
        "拖拽/点击上传",
        type=["pdf", "docx", "txt", "pptx", "ppt"],
        accept_multiple_files=True,
        label_visibility="collapsed",
        key=f"file_uploader_{st.session_state.uploader_key}"
    )
    st.caption("单文件最大200MB")

    # 3. 导入到知识库按钮（删除侧边栏的成功提示）
    if uploaded_files:
        if st.button("✅ 导入到知识库", use_container_width=True):
            with st.spinner("处理中..."):
                os.makedirs(DATA_DIR, exist_ok=True)
                vs = VectorStore()
                all_chunks = []
                
                for uploaded_file in uploaded_files:
                    save_path = os.path.join(DATA_DIR, uploaded_file.name)
                    with open(save_path, "wb") as f:
                        f.write(uploaded_file.getbuffer())
                    
                    temp_path = f"./temp_{uploaded_file.name}"
                    shutil.copyfile(save_path, temp_path)
                    
                    file_ext = temp_path.lower()
                    split_docs = []
                    
                    if file_ext.endswith(".pdf"):
                        loader = PyPDFLoader(temp_path)
                        split_docs = loader.load()
                    elif file_ext.endswith(".docx"):
                        loader = Docx2txtLoader(temp_path)
                        split_docs = loader.load()
                    elif file_ext.endswith(".txt"):
                        loader = TextLoader(temp_path)
                        split_docs = loader.load()
                    elif file_ext.endswith(".pptx") or file_ext.endswith(".ppt"):
                        native_ppt_docs = load_ppt_native(temp_path)
                        from langchain_core.documents import Document
                        split_docs = [Document(page_content=d["page_content"], metadata=d["metadata"]) for d in native_ppt_docs]
                    else:
                        st.error(f"不支持格式: {uploaded_file.name}")
                        os.remove(temp_path)
                        continue
                    
                    text_splitter = RecursiveCharacterTextSplitter(
                        chunk_size=600, chunk_overlap=150,
                        separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?", " ", ""]
                    )
                    split_docs = text_splitter.split_documents(split_docs)
                    
                    for doc in split_docs:
                        all_chunks.append({
                            "content": doc.page_content,
                            "metadata": {
                                "filename": uploaded_file.name,
                                "page": doc.metadata.get("page", doc.metadata.get("page_number", "unknown"))
                            }
                        })
                    
                    os.remove(temp_path)
                
                if all_chunks:
                    vs.add_documents(all_chunks)
                    st.session_state.rag_agent = RAGAgent(model=MODEL_NAME, top_k=TOP_K)
                    # 【修改】删除侧边栏的st.success，只设置标志位
                    st.session_state.show_import_success = True
                    st.rerun()
                else:
                    st.error("❌ 无有效内容")

    # 4. 清空知识库按钮
    if st.button("🧹 清空知识库", use_container_width=True):
        with st.spinner("清空中..."):
            vs = VectorStore()
            vs.clear_collection()
            
            if os.path.exists(DATA_DIR):
                for filename in os.listdir(DATA_DIR):
                    file_path = os.path.join(DATA_DIR, filename)
                    try:
                        if os.path.isfile(file_path) or os.path.islink(file_path):
                            os.unlink(file_path)
                        elif os.path.isdir(file_path):
                            shutil.rmtree(file_path)
                    except Exception as e:
                        st.error(f"删除失败：{filename}")
            
            st.session_state.rag_agent = RAGAgent(model=MODEL_NAME, top_k=TOP_K)
            st.session_state.uploader_key += 1
            st.success("✅ 知识库已清空")
            st.rerun()

    # 5. 功能导航区域
    st.divider()
    st.markdown('<h5 class="sidebar-title-nav">📋 功能导航</h5>', unsafe_allow_html=True)
    # 智能出题按钮
    if st.button("📝 智能出题", 
                 type="primary" if st.session_state.page_state == "main" else "secondary",
                 use_container_width=True):
        st.session_state.page_state = "main"
        st.rerun()
    # 智能问答按钮
    if st.button("💬 智能问答",
                 type="primary" if st.session_state.page_state == "qa" else "secondary",
                 use_container_width=True):
        st.session_state.page_state = "qa"
        st.rerun()
    # 历史试卷按钮
    if st.button("📜 历史试卷",
                 type="primary" if st.session_state.page_state == "history_exam" else "secondary",
                 use_container_width=True):
        st.session_state.page_state = "history_exam"
        st.rerun()
    # 知识点分析按钮
    if st.button("📊 知识点分析",
                 type="primary" if st.session_state.page_state == "knowledge_analysis" else "secondary",
                 use_container_width=True):
        st.session_state.page_state = "knowledge_analysis"
        st.rerun()
    # 我的错题本按钮
    if st.button("❌ 我的错题本",
                 type="primary" if st.session_state.page_state == "wrong_book" else "secondary",
                 use_container_width=True):
        st.session_state.page_state = "wrong_book"
        st.session_state.wrong_questions = [
            record for record in st.session_state.user_answer_records
            if not record["is_correct"]
        ]
        st.rerun()

# ========== 页面路由渲染 ==========
# 1. 智能出题主页面
if st.session_state.page_state == "main":
    # 【新增】检查并显示导入成功提示
    if st.session_state.show_import_success:
        st.success("✅ 导入成功，可以生成题目")
        # 显示一次后立即重置，避免刷新后重复显示
        st.session_state.show_import_success = False
    
    # 主标题
    st.markdown("<h1 style='text-align: center;'>📑 智能出题与在线测试</h1>", unsafe_allow_html=True)
    st.divider()

    # 横向出题配置面板
    col_knowledge, col_type, col_difficulty, col_num = st.columns([2.5, 1.5, 1.5, 1.5])
    with col_knowledge:
        knowledge_point = st.text_input(
            "核心知识点",
            placeholder="例如：唯物论、导数、Python",
            label_visibility="visible"
        )
    with col_type:
        q_type = st.selectbox("题型", ["单选", "多选", "判断", "简答"], index=0, label_visibility="visible")
    with col_difficulty:
        difficulty = st.selectbox("难度", ["简单", "中等", "进阶"], index=1, label_visibility="visible")
    with col_num:
        num_questions = st.slider("题目数量", min_value=1, max_value=20, value=5, label_visibility="visible")

    # 生成试题按钮
    generate_btn = st.button("🚀 生成试题", use_container_width=True, type="primary")
    st.divider()
    st.caption("🟣 智能题库系统 | 适配全学科学生使用 | 基于RAG与大语言模型")

    # 试题生成逻辑
    if generate_btn:
        doc_count = st.session_state.rag_agent.vector_store.get_collection_count()
        if doc_count == 0:
            st.error("❌ 知识库为空，请先上传课程材料并导入到知识库")
        else:
            with st.spinner(f"正在生成{num_questions}道试题..."):
                agent = st.session_state.rag_agent
                result = agent.generate_qa(
                    knowledge_point=knowledge_point,
                    q_type=q_type,
                    difficulty=difficulty,
                    num=num_questions
                )
                
                if "❌" in result:
                    st.error(result)
                else:
                    lines = [line.strip() for line in result.split("\n") if line.strip()]
                    st.session_state.generated_questions = lines
                    st.session_state.current_q_type = q_type
                    st.session_state.exam_mode = "view"
                    st.session_state.user_answers = {}
                    st.session_state.graded = False
                    st.session_state.score_detail = []
                    st.session_state.total_score = 0
                    exam_id = save_exam_history(lines, q_type, knowledge_point, difficulty, agent)
                    st.session_state.current_main_exam_id = exam_id
                    st.success(f"✅ 试题生成完成！已保存到历史试卷")

    # 试题展示区域
    st.divider()
    st.subheader("📝 生成的试题库")
    questions = st.session_state.generated_questions
    current_q_type = st.session_state.current_q_type
    total_question_num = len(questions)

    if questions:
        # 操作按钮
        col1, col2, _ = st.columns([1, 1, 3])
        with col1:
            if st.button("📝 在线测试", use_container_width=True, type="primary"):
                st.session_state.exam_mode = "test"
                st.session_state.graded = False
                st.session_state.user_answers = {}
                st.session_state.score_detail = []
                st.session_state.total_score = 0
                st.rerun()
        with col2:
            if st.button("📌 查看答案", use_container_width=True):
                st.session_state.exam_mode = "view"
                st.rerun()
        
        st.divider()
        
        # 查看答案模式
        if st.session_state.exam_mode == "view":
            render_view_mode(questions, current_q_type)
        # 在线测试模式
        elif st.session_state.exam_mode == "test":
            if not st.session_state.graded:
                submitted, updated_answers = render_test_mode(
                    questions=questions,
                    current_q_type=current_q_type,
                    total_question_num=total_question_num,
                    key_prefix="main",
                    user_answers=st.session_state.user_answers
                )
                st.session_state.user_answers = updated_answers
                
                if submitted:
                    total_score, score_detail = grade_exam_answers(
                        user_answers=st.session_state.user_answers,
                        current_q_type=current_q_type,
                        total_question_num=total_question_num,
                        model=st.session_state.model
                    )
                    st.session_state.total_score = total_score
                    st.session_state.score_detail = score_detail
                    st.session_state.graded = True
                    
                    main_test_record = {
                        "test_id": str(uuid.uuid4())[:8],
                        "test_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "total_score": total_score,
                        "score_detail": score_detail,
                        "user_answers": st.session_state.user_answers,
                        "q_type": current_q_type,
                        "questions": questions
                    }
                    
                    # ==================== 主界面测试 ====================
                    for detail in score_detail:
                        knowledge_points = extract_knowledge_points(detail["title"], detail["analysis"])
                        question_full_line = questions[detail["index"] - 1]
                        answer_record = {
                            "knowledge_points": knowledge_points,
                            "total_score": detail["total_score"],
                            "user_score": detail["get_score"],
                            "question_title": detail["title"],
                            "is_correct": detail["is_correct"],
                            "question_type": current_q_type,
                            "question_full_line": question_full_line 
                        }
                        st.session_state.user_answer_records.append(answer_record)
                    
                    if st.session_state.current_main_exam_id:
                        for idx, item in enumerate(st.session_state.exam_history):
                            if item["id"] == st.session_state.current_main_exam_id:
                                st.session_state.exam_history[idx]["main_test_record"] = main_test_record
                                break
                    
                    st.rerun()
            # 批改结果
            else:
                st.markdown('<a name="graded_results_main"></a>', unsafe_allow_html=True)
                st.divider()
                st.markdown(f"<h2 style='text-align: center;'>📊 测试结果：{st.session_state.total_score} / 100 分</h2>", unsafe_allow_html=True)
                score = st.session_state.total_score
                if score >= 90:
                    level = "优秀"
                    level_color = "#008000"
                elif score >= 80:
                    level = "良好"
                    level_color = "#32cd32"
                elif score >= 70:
                    level = "中等"
                    level_color = "#ffd700"
                elif score >= 60:
                    level = "及格"
                    level_color = "#ffa500"
                else:
                    level = "不及格"
                    level_color = "#ff4444"
                st.markdown(f"<h3 style='text-align: center; color: {level_color};'>等级：{level}</h3>", unsafe_allow_html=True)
                st.divider()
                st.subheader("📝 详细批改")
                for i in range(total_question_num):
                    if i not in st.session_state.user_answers:
                        continue
                    data = st.session_state.user_answers[i]
                    title = data["title"]
                    user_ans = data["user_ans"]
                    correct_ans = data["correct_ans"]
                    analysis = data["analysis"]
                    opts = data.get("opts", [])
                    detail = next((d for d in st.session_state.score_detail if d["index"] == i+1), None)
                    
                    with st.container(border=True):
                        st.markdown(f"**Q{i+1}. {title}**")
                        
                        if current_q_type in ["单选", "多选"] and len(opts) == 4:
                            st.markdown(f"A. {opts[0]}")
                            st.markdown(f"B. {opts[1]}")
                            st.markdown(f"C. {opts[2]}")
                            st.markdown(f"D. {opts[3]}")
                            st.markdown("<br>", unsafe_allow_html=True)
                        
                        if detail:
                            st.markdown(f"**本题总分：{detail['total_score']}分 | 你的得分：{detail['get_score']}分**")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            user_display = format_answer(user_ans)
                            st.markdown(f"<p style='color: #ff4444; font-weight: bold;'>你的答案：{user_display}</p>", unsafe_allow_html=True)
                        with col2:
                            st.markdown(f"<p style='color: #0066cc; font-weight: bold;'>正确答案：{correct_ans}</p>", unsafe_allow_html=True)
                        
                        if detail:
                            status = detail["status"]
                            if status == "full":
                                st.success(f"✅ 回答正确！{detail['remark']}")
                            elif status == "partial":
                                st.warning(f"⚠️ 部分正确 | {detail['remark']}")
                            else:
                                st.error(f"❌ 回答错误 | {detail['remark']}")
                        
                        if current_q_type == "简答" and detail and detail.get("hit_detail"):
                            with st.expander("🔍 查看采分点语义匹配详情", expanded=False):
                                st.markdown("**采分点命中详情：**")
                                for hd in detail["hit_detail"]:
                                    hit_status = "✅ 命中" if hd["is_hit"] else "❌ 未命中"
                                    color = "#008000" if hd["is_hit"] else "#ff4444"
                                    st.markdown(f"- 采分点：{hd['score_point']}")
                                    st.markdown(f"  匹配模式：{hd['mode']} | 相似度：<span style='color: {color}; font-weight: bold;'>{hd['similarity']}%</span> | 状态：{hit_status}", unsafe_allow_html=True)
                        
                        with st.expander("📌 查看详细解析", expanded=False):
                            formatted_analysis = analysis.replace('<br>', '\n').replace('\n', '<br>')
                            st.markdown(f"<p style='color: #008000;'><strong>详细解析：</strong><br>{formatted_analysis}</p>", unsafe_allow_html=True)
                    
                    st.markdown("<br>", unsafe_allow_html=True)
                st.divider()
                col_export, _ = st.columns([1, 1])
                with col_export:
                    export_clicked = st.button("📤 导出测试结果", use_container_width=True)
                    if export_clicked:
                        export_text = f"测试总分：{st.session_state.total_score}/100 分\n等级：{level}\n\n详细批改：\n"
                        for detail in st.session_state.score_detail:
                            export_text += f"Q{detail['index']}：总分{detail['total_score']}分，得分{detail['get_score']}分，{detail['remark']}\n"
                        st.download_button(
                            label="下载结果文件",
                            data=export_text,
                            file_name=f"考试结果_{st.session_state.total_score}分.txt",
                            mime="text/plain",
                            use_container_width=True
                        )
    else:
        st.info("暂无试题，请先导入课程材料并生成试题")

# 2. 智能问答页面
elif st.session_state.page_state == "qa":
    st.title("💬 课程智能问答助手")
    st.caption("基于你上传的课程材料，严格依据内容回答问题，支持PPT内容深度问答")
    st.divider()

    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])
    
    if prompt := st.chat_input("请输入你的问题..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        with st.chat_message("assistant"):
            with st.spinner("正在检索并生成回答..."):
                agent = st.session_state.rag_agent
                answer = agent.answer(prompt)
                st.markdown(answer)
                st.session_state.messages.append({"role": "assistant", "content": answer})
    
    st.divider()
    if st.button("🗑️ 清空全部问答记录", use_container_width=False):
        st.session_state.messages = []
        st.success("✅ 全部问答记录已清空")
        st.rerun()

# 3. 历史试卷页面
elif st.session_state.page_state == "history_exam":
    st.title("📝 历史试卷中心")
    st.divider()
    
    if st.session_state.exam_history:
        if st.button("🗑️ 清空全部试题记录", use_container_width=False):
            st.session_state.exam_history = []
            st.session_state.history_current_exam_id = ""
            st.session_state.history_questions = []
            st.session_state.history_q_type = ""
            st.success("✅ 全部试题记录已清空")
            st.rerun()
        st.divider()
    
    if st.session_state.history_current_exam_id:
        current_exam = next(
            (item for item in st.session_state.exam_history if item["id"] == st.session_state.history_current_exam_id),
            None
        )
        if current_exam:
            st.subheader(f"📝 试卷详情 | {current_exam['doc_title']} | {current_exam['q_type']}")
            if st.button("← 返回历史试卷列表", use_container_width=False):
                st.session_state.history_current_exam_id = ""
                st.session_state.history_questions = []
                st.session_state.history_q_type = ""
                st.rerun()
            st.divider()
            
            main_test_record = current_exam["main_test_record"]
            if main_test_record:
                st.subheader("📊 测试记录")
                st.metric(
                    label="测试得分",
                    value=f"{main_test_record['total_score']}/100分",
                    delta=f"测试时间：{main_test_record['test_time']}"
                )
                
                col_view, col_export = st.columns([1,1])
                with col_view:
                    if st.button("📋 查看答题详情", use_container_width=True):
                        st.session_state.history_questions = main_test_record["questions"]
                        st.session_state.history_q_type = main_test_record["q_type"]
                        st.session_state.history_user_answers = main_test_record["user_answers"]
                        st.session_state.history_total_score = main_test_record["total_score"]
                        st.session_state.history_score_detail = main_test_record["score_detail"]
                with col_export:
                    if st.button("📤 导出测试结果", use_container_width=True):
                        level = "优秀" if main_test_record['total_score'] >=90 else "良好" if main_test_record['total_score'] >=80 else "中等" if main_test_record['total_score'] >=70 else "及格" if main_test_record['total_score'] >=60 else "不及格"
                        export_text = f"试卷：{current_exam['doc_title']}\n测试时间：{main_test_record['test_time']}\n总分：{main_test_record['total_score']}/100 分\n等级：{level}\n\n详细批改：\n"
                        for detail in main_test_record["score_detail"]:
                            export_text += f"Q{detail['index']}：总分{detail['total_score']}分，得分{detail['get_score']}分，{detail['remark']}\n"
                        st.download_button(
                            label="下载结果文件",
                            data=export_text,
                            file_name=f"历史试卷测试结果_{main_test_record['total_score']}分.txt",
                            mime="text/plain",
                            use_container_width=True
                        )
                st.divider()

                if st.session_state.history_user_answers:
                    st.subheader("📝 答题详情")
                    total_num = len(st.session_state.history_questions)
                    for i in range(total_num):
                        if i not in st.session_state.history_user_answers:
                            continue
                        data = st.session_state.history_user_answers[i]
                        title = data["title"]
                        user_ans = data["user_ans"]
                        correct_ans = data["correct_ans"]
                        analysis = data["analysis"]
                        opts = data.get("opts", [])
                        detail = next((d for d in st.session_state.history_score_detail if d["index"] == i+1), None)
                        
                        with st.container(border=True):
                            st.markdown(f"**Q{i+1}. {title}**")
                            
                            if st.session_state.history_q_type in ["单选", "多选"] and len(opts) == 4:
                                st.markdown(f"A. {opts[0]}")
                                st.markdown(f"B. {opts[1]}")
                                st.markdown(f"C. {opts[2]}")
                                st.markdown(f"D. {opts[3]}")
                                st.markdown("<br>", unsafe_allow_html=True)
                            
                            if detail:
                                st.markdown(f"**本题总分：{detail['total_score']}分 | 你的得分：{detail['get_score']}分**")
                            
                            col1, col2 = st.columns(2)
                            with col1:
                                user_display = format_answer(user_ans)
                                st.markdown(f"<p style='color: #ff4444; font-weight: bold;'>你的答案：{user_display}</p>", unsafe_allow_html=True)
                            with col2:
                                st.markdown(f"<p style='color: #0066cc; font-weight: bold;'>正确答案：{correct_ans}</p>", unsafe_allow_html=True)
                            
                            if detail:
                                status = detail["status"]
                                if status == "full":
                                    st.success(f"✅ 回答正确！{detail['remark']}")
                                elif status == "partial":
                                    st.warning(f"⚠️ 部分正确 | {detail['remark']}")
                                else:
                                    st.error(f"❌ 回答错误 | {detail['remark']}")
                            
                            with st.expander("📌 查看详细解析", expanded=False):
                                formatted_analysis = analysis.replace('<br>', '\n').replace('\n', '<br>')
                                st.markdown(f"<p style='color: #008000;'><strong>详细解析：</strong><br>{formatted_analysis}</p>", unsafe_allow_html=True)
                        
                        st.markdown("<br>", unsafe_allow_html=True)
                st.divider()
            
            st.subheader("📌 题目与答案")
            render_view_mode(current_exam["questions"], current_exam["q_type"])
        else:
            st.error("❌ 未找到该试卷，请返回历史试卷列表")
            if st.button("← 返回历史试卷列表"):
                st.session_state.history_current_exam_id = ""
                st.rerun()
    else:
        if not st.session_state.exam_history:
            st.info("暂无试题历史记录，请先生成试题")
        else:
            for idx, item in enumerate(st.session_state.exam_history):
                exam_id = item["id"]
                q_type = item["q_type"]
                doc_title = item["doc_title"]
                create_time = item["time"]
                main_test_record = item["main_test_record"]
                is_tested = main_test_record is not None
                test_tag = "✅ 已测试" if is_tested else "⚠️ 未测试"
                
                panel_title = f"{test_tag} | 试卷 {exam_id} | {q_type} | {doc_title} | {create_time}"
                with st.expander(panel_title, expanded=False):
                    if st.button(f"❌ 删除本组", key=f"del_exam_{exam_id}"):
                        st.session_state.exam_history.pop(idx)
                        st.success("✅ 本组试题已删除")
                        st.rerun()
                    
                    st.divider()
                    
                    if st.button("📋 查看试卷详情", use_container_width=True, key=f"view_exam_{exam_id}"):
                        st.session_state.history_current_exam_id = exam_id
                        st.session_state.history_questions = item["questions"]
                        st.session_state.history_q_type = item["q_type"]
                        st.session_state.history_user_answers = {}
                        st.session_state.history_total_score = 0
                        st.session_state.history_score_detail = []
                        st.rerun()

# 4. 知识点分析页面
elif st.session_state.page_state == "knowledge_analysis":
    render_knowledge_analysis_page()

# 5. 我的错题本页面
elif st.session_state.page_state == "wrong_book":
    st.title("❌ 我的错题本")
    st.divider()

    wrong_questions = st.session_state.wrong_questions
    unique_wrongs = []
    seen_titles = set()
    for q in wrong_questions:
        if q["question_title"] not in seen_titles:
            seen_titles.add(q["question_title"])
            unique_wrongs.append(q)
    wrong_questions = unique_wrongs

    if not wrong_questions:
        st.info("🎉 太棒了！你目前没有错题，继续保持~")
    else:
        total_wrong = len(wrong_questions)
        knowledge_count = len(set([kp for q in wrong_questions for kp in q["knowledge_points"]]))
        col1, col2 = st.columns(2)
        with col1:
            st.metric(label="累计错题数", value=total_wrong)
        with col2:
            st.metric(label="涉及知识点数", value=knowledge_count)
        st.divider()

        st.subheader("📝 错题重测")
        st.caption("勾选需要重新测试的错题，点击「开始重测」即可作答")
        selected_ids = []
        with st.container(border=True):
            for idx, q in enumerate(wrong_questions):
                q_id = f"wrong_{idx}"
                kp_text = "、".join(q["knowledge_points"]) if q["knowledge_points"] else "未分类知识点"
                checked = st.checkbox(
                    f"Q{idx+1} | {q['question_title']}（知识点：{kp_text}）",
                    key=q_id,
                    value=True
                )
                if checked:
                    selected_ids.append(idx)
        
        col_start, col_clear = st.columns([1,1])
        with col_start:
            if st.button("🚀 开始重测", use_container_width=True, type="primary"):
                if not selected_ids:
                    st.warning("请至少选择一道错题")
                else:
                    selected_questions = [wrong_questions[i] for i in selected_ids]
                    st.session_state.wrong_questions_for_test = selected_questions
                    st.session_state.wrong_exam_mode = "test"
                    st.session_state.wrong_user_answers = {}
                    st.session_state.wrong_graded = False
                    st.session_state.wrong_total_score = 0
                    st.session_state.wrong_score_detail = []
                    st.rerun()
        with col_clear:
            if st.button("🗑️ 清空全部错题", use_container_width=True):
                st.session_state.user_answer_records = [
                    record for record in st.session_state.user_answer_records
                    if record["is_correct"]
                ]
                st.session_state.wrong_questions = []
                st.success("✅ 全部错题已清空")
                st.rerun()
        st.divider()

        if st.session_state.wrong_exam_mode == "test":
            st.subheader("📝 错题重测作答")
            test_questions = st.session_state.wrong_questions_for_test
            total_num = len(test_questions)
            q_type = test_questions[0]["question_type"] if test_questions else "单选"

            # ==================== 错题本====================
            submitted, updated_answers = render_test_mode(
                # 修改：传入 question_full_line 而不是 question_title
                questions=[q["question_full_line"] for q in test_questions],
                current_q_type=q_type,
                total_question_num=total_num,
                key_prefix="wrong",
                user_answers=st.session_state.wrong_user_answers
            )
            
            st.session_state.wrong_user_answers = updated_answers

            if submitted:
                total_score, score_detail = grade_exam_answers(
                    user_answers=st.session_state.wrong_user_answers,
                    current_q_type=q_type,
                    total_question_num=total_num,
                    model=st.session_state.model
                )
                st.session_state.wrong_total_score = total_score
                st.session_state.wrong_score_detail = score_detail
                st.session_state.wrong_graded = True
                st.session_state.wrong_exam_mode = "graded"
                st.rerun()

        elif st.session_state.wrong_exam_mode == "graded":
            st.markdown(f"<h2 style='text-align: center;'>📊 重测结果：{st.session_state.wrong_total_score} / 100 分</h2>", unsafe_allow_html=True)
            score = st.session_state.wrong_total_score
            if score >= 90:
                level = "优秀"
                level_color = "#008000"
            elif score >= 80:
                level = "良好"
                level_color = "#32cd32"
            elif score >= 70:
                level = "中等"
                level_color = "#ffd700"
            elif score >= 60:
                level = "及格"
                level_color = "#ffa500"
            else:
                level = "不及格"
                level_color = "#ff4444"
            st.markdown(f"<h3 style='text-align: center; color: {level_color};'>等级：{level}</h3>", unsafe_allow_html=True)
            st.divider()

            st.subheader("📝 详细批改")
            test_questions = st.session_state.wrong_questions_for_test
            total_num = len(test_questions)
            q_type = test_questions[0]["question_type"] if test_questions else "单选"

            for i in range(total_num):
                if i not in st.session_state.wrong_user_answers:
                    continue
                data = st.session_state.wrong_user_answers[i]
                title = data["title"]
                user_ans = data["user_ans"]
                correct_ans = data["correct_ans"]
                analysis = data["analysis"]
                opts = data.get("opts", [])
                detail = next((d for d in st.session_state.wrong_score_detail if d["index"] == i+1), None)
                
                with st.container(border=True):
                    st.markdown(f"**Q{i+1}. {title}**")
                    
                    if q_type in ["单选", "多选"] and len(opts) == 4:
                        st.markdown(f"A. {opts[0]}")
                        st.markdown(f"B. {opts[1]}")
                        st.markdown(f"C. {opts[2]}")
                        st.markdown(f"D. {opts[3]}")
                        st.markdown("<br>", unsafe_allow_html=True)
                    
                    if detail:
                        st.markdown(f"**本题总分：{detail['total_score']}分 | 你的得分：{detail['get_score']}分**")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        user_display = format_answer(user_ans)
                        st.markdown(f"<p style='color: #ff4444; font-weight: bold;'>你的答案：{user_display}</p>", unsafe_allow_html=True)
                    with col2:
                        st.markdown(f"<p style='color: #0066cc; font-weight: bold;'>正确答案：{correct_ans}</p>", unsafe_allow_html=True)
                    
                    if detail:
                        status = detail["status"]
                        if status == "full":
                            st.success(f"✅ 回答正确！{detail['remark']}")
                        elif status == "partial":
                            st.warning(f"⚠️ 部分正确 | {detail['remark']}")
                        else:
                            st.error(f"❌ 回答错误 | {detail['remark']}")
                    
                    with st.expander("📌 查看详细解析", expanded=False):
                        formatted_analysis = analysis.replace('<br>', '\n').replace('\n', '<br>')
                        st.markdown(f"<p style='color: #008000;'><strong>详细解析：</strong><br>{formatted_analysis}</p>", unsafe_allow_html=True)
                
                st.markdown("<br>", unsafe_allow_html=True)
            
            st.divider()
            col_retest, col_export, col_back = st.columns([1,1,1])
            with col_retest:
                if st.button("🔄 重新测试", use_container_width=True, type="primary"):
                    st.session_state.wrong_exam_mode = "test"
                    st.session_state.wrong_user_answers = {}
                    st.session_state.wrong_graded = False
                    st.session_state.wrong_total_score = 0
                    st.session_state.wrong_score_detail = []
                    st.rerun()
            with col_export:
                if st.button("📤 导出重测结果", use_container_width=True):
                    export_text = f"错题重测结果\n总分：{st.session_state.wrong_total_score}/100 分\n等级：{level}\n\n详细批改：\n"
                    for detail in st.session_state.wrong_score_detail:
                        export_text += f"Q{detail['index']}：总分{detail['total_score']}分，得分{detail['get_score']}分，{detail['remark']}\n"
                    st.download_button(
                        label="下载结果文件",
                        data=export_text,
                        file_name=f"错题重测结果_{st.session_state.wrong_total_score}分.txt",
                        mime="text/plain",
                        use_container_width=True
                    )
            with col_back:
                if st.button("← 返回错题列表", use_container_width=True):
                    st.session_state.wrong_exam_mode = "view"
                    st.session_state.wrong_user_answers = {}
                    st.session_state.wrong_graded = False
                    st.rerun()

        else:
            st.subheader("📋 错题详情列表")
            for idx, q in enumerate(wrong_questions):
                kp_text = "、".join(q["knowledge_points"]) if q["knowledge_points"] else "未分类知识点"
                with st.expander(f"Q{idx+1} | 知识点：{kp_text}", expanded=False):
                    st.markdown(f"**题干：{q['question_title']}**")
                    st.markdown(f"**题型：{q['question_type']}**")
                    st.markdown(f"**你的得分：{q['user_score']}/{q['total_score']}分**")
                    st.markdown(f"**是否答对：❌ 答错**")
                    st.divider()
                    st.markdown("**💡 知识点巩固建议**")
                    st.info(f"本题核心考察「{kp_text}」，建议重新回顾对应知识点的课程内容，结合解析理解考点，再通过重测巩固掌握。")

# 全局底部版权信息
st.divider()
st.caption("💡 NLP智能题库系统 | 基于Sentence-BERT语义匹配 | 在线作答 | 精准评分 | 自动批改")
'''
# ========== 最顶部：强制国内镜像配置 ==========
import os
os.environ['HF_ENDPOINT'] = 'https://hf-mirror.com'
os.environ['HF_HOME'] = './hf_cache'

import streamlit as st
import shutil
import re
import uuid
from datetime import datetime
from rag_agent import RAGAgent
from config import MODEL_NAME, TOP_K, DATA_DIR
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from pptx import Presentation
from vector_store import VectorStore

# ========== 导入拆分的模块 ==========
from exam_core import clean_option_prefix, grade_exam_answers, format_answer
from exam_ui import render_view_mode, render_test_mode, render_graded_results

# ========== 知识点分析模块 ==========
from knowledge_analysis import render_knowledge_analysis_page, extract_knowledge_points

# ========== 阿里云向量 API（替代本地 Sentence-BERT）==========
import dashscope
from dashscope import TextEmbedding

# ========== 全局SSL配置 ==========
import ssl
ssl._create_default_https_context = ssl._create_unverified_context

# ========== 阿里云密钥配置 ==========
dashscope.api_key = "sk-737ae41763b4442dadcb9f6b1de03e4c"

# ========== 阿里云向量 API 封装 ==========
def get_ali_embedding(text: str):
    try:
        resp = TextEmbedding.call(model="text-embedding-v2", input=text)
        if resp.status_code == 200:
            return resp.output["embeddings"][0]["embedding"]
        else:
            return []
    except:
        return []

# ========== 工具函数 ==========
def save_exam_history(questions: list, q_type: str, knowledge_point: str, difficulty: str, agent):
    if "exam_history" not in st.session_state:
        st.session_state.exam_history = []
    
    if knowledge_point.strip():
        doc_title = knowledge_point
    else:
        try:
            doc_names = agent.vector_store.get_all_filenames()
            doc_title = "、".join(doc_names) if doc_names else "全部知识点"
        except:
            doc_title = "全部知识点"
    
    exam_id = str(uuid.uuid4())[:8]
    history_item = {
        "id": exam_id,
        "time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "doc_title": doc_title,
        "knowledge_point": knowledge_point,
        "q_type": q_type,
        "difficulty": difficulty,
        "questions": questions,
        "main_test_record": None,
        "practice_records": []
    }
    st.session_state.exam_history.insert(0, history_item)
    return exam_id

def load_ppt_native(file_path: str):
    prs = Presentation(file_path)
    docs = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        text_content = "\n".join(slide_text)
        if text_content.strip():
            docs.append({
                "page_content": text_content,
                "metadata": {"page": slide_num, "filename": os.path.basename(file_path)}
            })
    return docs

# ========== 页面配置 ==========
st.set_page_config(
    page_title="智能题库系统",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ========== 【舒适版】全局布局CSS ==========
st.markdown("""
<style>
/* ========== 1. 彻底隐藏侧边栏折叠按钮，消除顶部空隙 ========== */
[data-testid="stSidebarCollapseButton"] {
    display: none !important;
}
[data-testid="stSidebarHeader"] {
    display: none !important;
    height: 0 !important;
    padding: 0 !important;
    margin: 0 !important;
}

/* ========== 2. 侧边栏容器：舒适的边距 ========== */
[data-testid="stSidebar"] {
    padding: 0 !important;
    margin: 0 !important;
}
[data-testid="stSidebar"] > div:first-child {
    padding-top: 0.8rem !important;
    padding-bottom: 0.8rem !important;
    padding-left: 1rem !important;
    padding-right: 1rem !important;
    margin: 0 !important;
    width: 100% !important;
    top: 0 !important;
    position: absolute !important;
    height: 100vh !important;
}

/* ========== 3. 侧边栏所有元素间距：舒适不拥挤 ========== */
[data-testid="stSidebar"] [data-testid="stVerticalBlock"] > div {
    gap: 0.4rem !important;
    padding: 0 !important;
    margin: 0 !important;
}

/* ========== 4. 侧边栏标题通用样式 ========== */
.sidebar-title-main {
    margin: 0 0 0.2rem 0 !important;
    padding: 0 !important;
    line-height: 1.3 !important;
}
.sidebar-title-sub {
    margin: 0 0 0.15rem 0 !important;
    padding: 0 !important;
    line-height: 1.3 !important;
}
.sidebar-title-nav {
    margin: 0 0 0.6rem 0 !important;
    padding: 0 !important;
    line-height: 1.3 !important;
}

/* ========== 5. 侧边栏所有标题通用margin控制 ========== */
[data-testid="stSidebar"] h1, 
[data-testid="stSidebar"] h2, 
[data-testid="stSidebar"] h3,
[data-testid="stSidebar"] h4,
[data-testid="stSidebar"] h5,
[data-testid="stSidebar"] h6 {
    margin-top: 0 !important;
    margin-bottom: 0.15rem !important;
    padding: 0 !important;
    line-height: 1.3 !important;
}

/* ========== 6. 按钮：舒适的高度和边距 ========== */
[data-testid="stSidebar"] button {
    margin-top: 0.1rem !important;
    margin-bottom: 0.1rem !important;
    min-height: 2.2rem !important;
    padding-top: 0.35rem !important;
    padding-bottom: 0.35rem !important;
    line-height: 1.2 !important;
}

/* ========== 7. 文件上传组件：舒适的边距 + 中文替换 ========== */
[data-testid="stSidebar"] [data-testid="stFileUploader"] {
    margin-top: 0.15rem !important;
    margin-bottom: 0.15rem !important;
    padding: 0 !important;
}
[data-testid="stSidebar"] [data-testid="stFileUploader"] > section {
    padding: 0.6rem 0.5rem !important;
    min-height: auto !important;
    margin: 0 !important;
}

/* ========== 8. 已上传文件列表：舒适的边距 ========== */
[data-testid="stSidebar"] [data-testid="stFileUploaderFile"] {
    margin-top: 0.1rem !important;
    margin-bottom: 0.1rem !important;
    padding-top: 0.2rem !important;
    padding-bottom: 0.2rem !important;
    min-height: auto !important;
}

/* ========== 9. 小字说明：舒适的边距 ========== */
[data-testid="stSidebar"] small {
    margin-top: 0.05rem !important;
    margin-bottom: 0.05rem !important;
    line-height: 1.2 !important;
}

/* ========== 10. 分割线：舒适的边距 ========== */
[data-testid="stSidebar"] hr {
    margin-top: 0.3rem !important;
    margin-bottom: 0.3rem !important;
    padding: 0 !important;
}

/* ========== 11. 文件上传框替换为中文 ========== */
[data-testid="stFileUploader"] section > div:first-child > span:first-child {
    font-size: 0 !important;
}
[data-testid="stFileUploader"] section > div:first-child > span:first-child::before {
    content: "拖拽文件到此处";
    font-size: 2.4rem;
    color: #262730;
}

[data-testid="stFileUploader"] section > div:first-child > small {
    font-size: 0 !important;
}
[data-testid="stFileUploader"] section > div:first-child > small::before {
    content: "单文件最大200MB • 支持 PDF, DOCX, TXT, PPT 等格式";
    font-size: 1.1rem;
    color: #6b7280;
}

[data-testid="stFileUploader"] section button p {
    font-size: 0 !important;
}
[data-testid="stFileUploader"] section button p::before {
    content: "选择文件";
    font-size: 1.4rem;
}

@media (prefers-color-scheme: dark) {
    [data-testid="stFileUploader"] section > div:first-child > span:first-child::before {
        color: #fafafa;
    }
    [data-testid="stFileUploader"] section > div:first-child > small::before {
        color: #9ca3af;
    }
}

/* ========== 主页面正确优化 ========== */
[data-testid="stMainBlockContainer"] {
    padding-top: 2rem !important;
    padding-bottom: 2rem !important;
    max-width: 90% !important;
}
[data-testid="stMainBlockContainer"] h1 {
    margin-top: 0.3rem !important;
    margin-bottom: 0.8rem !important;
    line-height: 1.3 !important;
}
[data-testid="stMainBlockContainer"] h2, 
[data-testid="stMainBlockContainer"] h3 {
    margin-top: 0.7rem !important;
    margin-bottom: 0.4rem !important;
    line-height: 1.3 !important;
}
[data-testid="stMainBlockContainer"] hr {
    margin-top: 0.7rem !important;
    margin-bottom: 0.7rem !important;
}
[data-testid="stMainBlockContainer"] [data-testid="stTextInput"],
[data-testid="stMainBlockContainer"] [data-testid="stSelectbox"],
[data-testid="stMainBlockContainer"] [data-testid="stSlider"] {
    margin-top: 0.05rem !important;
    margin-bottom: 0.05rem !important;
}
[data-testid="stMainBlockContainer"] button {
    margin-top: 0.1rem !important;
    margin-bottom: 0.1rem !important;
}
</style>
""", unsafe_allow_html=True)

# ========== session_state 全量初始化 ==========
if "initialized" not in st.session_state:
    st.cache_data.clear()
    st.cache_resource.clear()
    st.session_state.initialized = True
    st.session_state.page_state = "main"
    st.session_state.exam_history = []
    st.session_state.generated_questions = []
    st.session_state.current_q_type = ""
    st.session_state.exam_mode = "view"
    st.session_state.user_answers = {}
    st.session_state.graded = False
    st.session_state.score_detail = []
    st.session_state.total_score = 0
    if "user_answer_records" not in st.session_state:
        st.session_state.user_answer_records = []
    st.session_state.history_current_exam_id = ""
    st.session_state.history_questions = []
    st.session_state.history_q_type = ""
    if "messages" not in st.session_state:
        st.session_state.messages = []
    st.session_state.wrong_questions = []
    st.session_state.wrong_exam_mode = "view"
    st.session_state.wrong_user_answers = {}
    st.session_state.wrong_graded = False
    st.session_state.wrong_total_score = 0
    st.session_state.wrong_score_detail = []
    st.session_state.selected_wrong_ids = []
    if "uploader_key" not in st.session_state:
        st.session_state.uploader_key = 0
    if "show_import_success" not in st.session_state:
        st.session_state.show_import_success = False

# RAG Agent 初始化
if "rag_agent" not in st.session_state:
    st.session_state.rag_agent = RAGAgent(model=MODEL_NAME, top_k=TOP_K)

# ========== 【舒适版】左侧侧边栏 ==========
with st.sidebar:
    st.markdown('<h4 class="sidebar-title-main">🧠 智能题库系统</h4>', unsafe_allow_html=True)
    st.divider()

    st.markdown('<h5 class="sidebar-title-sub">📂 文档导入</h5>', unsafe_allow_html=True)
    st.caption("支持 PDF/DOCX/TXT/PPT/PPTX")
    uploaded_files = st.file_uploader(
        "拖拽/点击上传",
        type=["pdf", "docx", "txt", "pptx", "ppt"],
        accept_multiple_files=True,
        label_visibility="collapsed",
        key=f"file_uploader_{st.session_state.uploader_key}"
    )
    st.caption("单文件最大200MB")

    if uploaded_files:
        if st.button("✅ 导入到知识库", use_container_width=True):
            with st.spinner("处理中..."):
                os.makedirs(DATA_DIR, exist_ok=True)
                vs = VectorStore()
                all_chunks = []
                
                for uploaded_file in uploaded_files:
                    save_path = os.path.join(DATA_DIR, uploaded_file.name)
                    with open(save_path, "wb") as f:
                        f.write(uploaded_file.getbuffer())
                    
                    temp_path = f"./temp_{uploaded_file.name}"
                    shutil.copyfile(save_path, temp_path)
                    
                    file_ext = temp_path.lower()
                    split_docs = []
                    
                    if file_ext.endswith(".pdf"):
                        loader = PyPDFLoader(temp_path)
                        split_docs = loader.load()
                    elif file_ext.endswith(".docx"):
                        loader = Docx2txtLoader(temp_path)
                        split_docs = loader.load()
                    elif file_ext.endswith(".txt"):
                        loader = TextLoader(temp_path)
                        split_docs = loader.load()
                    elif file_ext.endswith(".pptx") or file_ext.endswith(".ppt"):
                        native_ppt_docs = load_ppt_native(temp_path)
                        from langchain_core.documents import Document
                        split_docs = [Document(page_content=d["page_content"], metadata=d["metadata"]) for d in native_ppt_docs]
                    else:
                        st.error(f"不支持格式: {uploaded_file.name}")
                        os.remove(temp_path)
                        continue
                    
                    text_splitter = RecursiveCharacterTextSplitter(
                        chunk_size=600, chunk_overlap=150,
                        separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?", " ", ""]
                    )
                    split_docs = text_splitter.split_documents(split_docs)
                    
                    for doc in split_docs:
                        all_chunks.append({
                            "content": doc.page_content,
                            "metadata": {
                                "filename": uploaded_file.name,
                                "page": doc.metadata.get("page", doc.metadata.get("page_number", "unknown"))
                            }
                        })
                    
                    os.remove(temp_path)
                
                if all_chunks:
                    vs.add_documents(all_chunks)
                    st.session_state.rag_agent = RAGAgent(model=MODEL_NAME, top_k=TOP_K)
                    st.session_state.show_import_success = True
                    st.rerun()
                else:
                    st.error("❌ 无有效内容")

    if st.button("🧹 清空知识库", use_container_width=True):
        with st.spinner("清空中..."):
            vs = VectorStore()
            vs.clear_collection()
            
            if os.path.exists(DATA_DIR):
                for filename in os.listdir(DATA_DIR):
                    file_path = os.path.join(DATA_DIR, filename)
                    try:
                        if os.path.isfile(file_path) or os.path.islink(file_path):
                            os.unlink(file_path)
                        elif os.path.isdir(file_path):
                            shutil.rmtree(file_path)
                    except Exception as e:
                        st.error(f"删除失败：{filename}")
            
            st.session_state.rag_agent = RAGAgent(model=MODEL_NAME, top_k=TOP_K)
            st.session_state.uploader_key += 1
            st.success("✅ 知识库已清空")
            st.rerun()

    st.divider()
    st.markdown('<h5 class="sidebar-title-nav">📋 功能导航</h5>', unsafe_allow_html=True)
    if st.button("📝 智能出题", 
                 type="primary" if st.session_state.page_state == "main" else "secondary",
                 use_container_width=True):
        st.session_state.page_state = "main"
        st.rerun()
    if st.button("💬 智能问答",
                 type="primary" if st.session_state.page_state == "qa" else "secondary",
                 use_container_width=True):
        st.session_state.page_state = "qa"
        st.rerun()
    if st.button("📜 历史试卷",
                 type="primary" if st.session_state.page_state == "history_exam" else "secondary",
                 use_container_width=True):
        st.session_state.page_state = "history_exam"
        st.rerun()
    if st.button("📊 知识点分析",
                 type="primary" if st.session_state.page_state == "knowledge_analysis" else "secondary",
                 use_container_width=True):
        st.session_state.page_state = "knowledge_analysis"
        st.rerun()
    if st.button("❌ 我的错题本",
                 type="primary" if st.session_state.page_state == "wrong_book" else "secondary",
                 use_container_width=True):
        st.session_state.page_state = "wrong_book"
        st.session_state.wrong_questions = [
            record for record in st.session_state.user_answer_records
            if not record["is_correct"]
        ]
        st.rerun()

# ========== 页面路由渲染 ==========
# 1. 智能出题主页面
if st.session_state.page_state == "main":
    if st.session_state.show_import_success:
        st.success("✅ 导入成功，可以生成题目")
        st.session_state.show_import_success = False
    
    st.markdown("<h1 style='text-align: center;'>📑 智能出题与在线测试</h1>", unsafe_allow_html=True)
    st.divider()

    col_knowledge, col_type, col_difficulty, col_num = st.columns([2.5, 1.5, 1.5, 1.5])
    with col_knowledge:
        knowledge_point = st.text_input(
            "核心知识点",
            placeholder="例如：唯物论、导数、Python",
            label_visibility="visible"
        )
    with col_type:
        q_type = st.selectbox("题型", ["单选", "多选", "判断", "简答"], index=0, label_visibility="visible")
    with col_difficulty:
        difficulty = st.selectbox("难度", ["简单", "中等", "进阶"], index=1, label_visibility="visible")
    with col_num:
        num_questions = st.slider("题目数量", min_value=1, max_value=20, value=5, label_visibility="visible")

    generate_btn = st.button("🚀 生成试题", use_container_width=True, type="primary")
    st.divider()
    st.caption("🟣 智能题库系统 | 适配全学科学生使用 | 基于RAG与大语言模型")

    if generate_btn:
        doc_count = st.session_state.rag_agent.vector_store.get_collection_count()
        if doc_count == 0:
            st.error("❌ 知识库为空，请先上传课程材料并导入到知识库")
        else:
            with st.spinner(f"正在生成{num_questions}道试题..."):
                agent = st.session_state.rag_agent
                result = agent.generate_qa(
                    knowledge_point=knowledge_point,
                    q_type=q_type,
                    difficulty=difficulty,
                    num=num_questions
                )
                
                if "❌" in result:
                    st.error(result)
                else:
                    lines = [line.strip() for line in result.split("\n") if line.strip()]
                    st.session_state.generated_questions = lines
                    st.session_state.current_q_type = q_type
                    st.session_state.exam_mode = "view"
                    st.session_state.user_answers = {}
                    st.session_state.graded = False
                    st.session_state.score_detail = []
                    st.session_state.total_score = 0
                    exam_id = save_exam_history(lines, q_type, knowledge_point, difficulty, agent)
                    st.session_state.current_main_exam_id = exam_id
                    st.success(f"✅ 试题生成完成！已保存到历史试卷")

    st.divider()
    st.subheader("📝 生成的试题库")
    questions = st.session_state.generated_questions
    current_q_type = st.session_state.current_q_type
    total_question_num = len(questions)

    if questions:
        col1, col2, _ = st.columns([1, 1, 3])
        with col1:
            if st.button("📝 在线测试", use_container_width=True, type="primary"):
                st.session_state.exam_mode = "test"
                st.session_state.graded = False
                st.session_state.user_answers = {}
                st.session_state.score_detail = []
                st.session_state.total_score = 0
                st.rerun()
        with col2:
            if st.button("📌 查看答案", use_container_width=True):
                st.session_state.exam_mode = "view"
                st.rerun()
        
        st.divider()
        
        if st.session_state.exam_mode == "view":
            render_view_mode(questions, current_q_type)
        elif st.session_state.exam_mode == "test":
            if not st.session_state.graded:
                submitted, updated_answers = render_test_mode(
                    questions=questions,
                    current_q_type=current_q_type,
                    total_question_num=total_question_num,
                    key_prefix="main",
                    user_answers=st.session_state.user_answers
                )
                st.session_state.user_answers = updated_answers
                
                if submitted:
                    total_score, score_detail = grade_exam_answers(
                        user_answers=st.session_state.user_answers,
                        current_q_type=current_q_type,
                        total_question_num=total_question_num,
                        model=get_ali_embedding
                    )
                    st.session_state.total_score = total_score
                    st.session_state.score_detail = score_detail
                    st.session_state.graded = True
                    
                    main_test_record = {
                        "test_id": str(uuid.uuid4())[:8],
                        "test_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "total_score": total_score,
                        "score_detail": score_detail,
                        "user_answers": st.session_state.user_answers,
                        "q_type": current_q_type,
                        "questions": questions
                    }
                    
                    for detail in score_detail:
                        knowledge_points = extract_knowledge_points(detail["title"], detail["analysis"])
                        question_full_line = questions[detail["index"] - 1]
                        answer_record = {
                            "knowledge_points": knowledge_points,
                            "total_score": detail["total_score"],
                            "user_score": detail["get_score"],
                            "question_title": detail["title"],
                            "is_correct": detail["is_correct"],
                            "question_type": current_q_type,
                            "question_full_line": question_full_line 
                        }
                        st.session_state.user_answer_records.append(answer_record)
                    
                    if st.session_state.current_main_exam_id:
                        for idx, item in enumerate(st.session_state.exam_history):
                            if item["id"] == st.session_state.current_main_exam_id:
                                st.session_state.exam_history[idx]["main_test_record"] = main_test_record
                                break
                    
                    st.rerun()
            else:
                st.markdown('<a name="graded_results_main"></a>', unsafe_allow_html=True)
                st.divider()
                st.markdown(f"<h2 style='text-align: center;'>📊 测试结果：{st.session_state.total_score} / 100 分</h2>", unsafe_allow_html=True)
                score = st.session_state.total_score
                if score >= 90:
                    level = "优秀"
                    level_color = "#008000"
                elif score >= 80:
                    level = "良好"
                    level_color = "#32cd32"
                elif score >= 70:
                    level = "中等"
                    level_color = "#ffd700"
                elif score >= 60:
                    level = "及格"
                    level_color = "#ffa500"
                else:
                    level = "不及格"
                    level_color = "#ff4444"
                st.markdown(f"<h3 style='text-align: center; color: {level_color};'>等级：{level}</h3>", unsafe_allow_html=True)
                st.divider()
                st.subheader("📝 详细批改")
                for i in range(total_question_num):
                    if i not in st.session_state.user_answers:
                        continue
                    data = st.session_state.user_answers[i]
                    title = data["title"]
                    user_ans = data["user_ans"]
                    correct_ans = data["correct_ans"]
                    analysis = data["analysis"]
                    opts = data.get("opts", [])
                    detail = next((d for d in st.session_state.score_detail if d["index"] == i+1), None)
                    
                    with st.container(border=True):
                        st.markdown(f"**Q{i+1}. {title}**")
                        
                        if current_q_type in ["单选", "多选"] and len(opts) == 4:
                            st.markdown(f"A. {opts[0]}")
                            st.markdown(f"B. {opts[1]}")
                            st.markdown(f"C. {opts[2]}")
                            st.markdown(f"D. {opts[3]}")
                            st.markdown("<br>", unsafe_allow_html=True)
                        
                        if detail:
                            st.markdown(f"**本题总分：{detail['total_score']}分 | 你的得分：{detail['get_score']}分**")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            user_display = format_answer(user_ans)
                            st.markdown(f"<p style='color: #ff4444; font-weight: bold;'>你的答案：{user_display}</p>", unsafe_allow_html=True)
                        with col2:
                            st.markdown(f"<p style='color: #0066cc; font-weight: bold;'>正确答案：{correct_ans}</p>", unsafe_allow_html=True)
                        
                        if detail:
                            status = detail["status"]
                            if status == "full":
                                st.success(f"✅ 回答正确！{detail['remark']}")
                            elif status == "partial":
                                st.warning(f"⚠️ 部分正确 | {detail['remark']}")
                            else:
                                st.error(f"❌ 回答错误 | {detail['remark']}")
                        
                        if current_q_type == "简答" and detail and detail.get("hit_detail"):
                            with st.expander("🔍 查看采分点语义匹配详情", expanded=False):
                                st.markdown("**采分点命中详情：**")
                                for hd in detail["hit_detail"]:
                                    hit_status = "✅ 命中" if hd["is_hit"] else "❌ 未命中"
                                    color = "#008000" if hd["is_hit"] else "#ff4444"
                                    st.markdown(f"- 采分点：{hd['score_point']}")
                                    st.markdown(f"  匹配模式：{hd['mode']} | 相似度：<span style='color: {color}; font-weight: bold;'>{hd['similarity']}%</span> | 状态：{hit_status}", unsafe_allow_html=True)
                        
                        with st.expander("📌 查看详细解析", expanded=False):
                            formatted_analysis = analysis.replace('<br>', '\n').replace('\n', '<br>')
                            st.markdown(f"<p style='color: #008000;'><strong>详细解析：</strong><br>{formatted_analysis}</p>", unsafe_allow_html=True)
                    
                    st.markdown("<br>", unsafe_allow_html=True)
                st.divider()
                col_export, _ = st.columns([1, 1])
                with col_export:
                    export_clicked = st.button("📤 导出测试结果", use_container_width=True)
                    if export_clicked:
                        export_text = f"测试总分：{st.session_state.total_score}/100 分\n等级：{level}\n\n详细批改：\n"
                        for detail in st.session_state.score_detail:
                            export_text += f"Q{detail['index']}：总分{detail['total_score']}分，得分{detail['get_score']}分，{detail['remark']}\n"
                        st.download_button(
                            label="下载结果文件",
                            data=export_text,
                            file_name=f"考试结果_{st.session_state.total_score}分.txt",
                            mime="text/plain",
                            use_container_width=True
                        )
    else:
        st.info("暂无试题，请先导入课程材料并生成试题")

# 2. 智能问答页面
elif st.session_state.page_state == "qa":
    st.title("💬 课程智能问答助手")
    st.caption("基于你上传的课程材料，严格依据内容回答问题，支持PPT内容深度问答")
    st.divider()

    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])
    
    if prompt := st.chat_input("请输入你的问题..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        with st.chat_message("assistant"):
            with st.spinner("正在检索并生成回答..."):
                agent = st.session_state.rag_agent
                answer = agent.answer(prompt)
                st.markdown(answer)
                st.session_state.messages.append({"role": "assistant", "content": answer})
    
    st.divider()
    if st.button("🗑️ 清空全部问答记录", use_container_width=False):
        st.session_state.messages = []
        st.success("✅ 全部问答记录已清空")
        st.rerun()

# 3. 历史试卷页面
elif st.session_state.page_state == "history_exam":
    st.title("📝 历史试卷中心")
    st.divider()
    
    if st.session_state.exam_history:
        if st.button("🗑️ 清空全部试题记录", use_container_width=False):
            st.session_state.exam_history = []
            st.session_state.history_current_exam_id = ""
            st.session_state.history_questions = []
            st.session_state.history_q_type = ""
            st.success("✅ 全部试题记录已清空")
            st.rerun()
        st.divider()
    
    if st.session_state.history_current_exam_id:
        current_exam = next(
            (item for item in st.session_state.exam_history if item["id"] == st.session_state.history_current_exam_id),
            None
        )
        if current_exam:
            st.subheader(f"📝 试卷详情 | {current_exam['doc_title']} | {current_exam['q_type']}")
            if st.button("← 返回历史试卷列表", use_container_width=False):
                st.session_state.history_current_exam_id = ""
                st.session_state.history_questions = []
                st.session_state.history_q_type = ""
                st.rerun()
            st.divider()
            
            main_test_record = current_exam["main_test_record"]
            if main_test_record:
                st.subheader("📊 测试记录")
                st.metric(
                    label="测试得分",
                    value=f"{main_test_record['total_score']}/100分",
                    delta=f"测试时间：{main_test_record['test_time']}"
                )
                
                col_view, col_export = st.columns([1,1])
                with col_view:
                    if st.button("📋 查看答题详情", use_container_width=True):
                        st.session_state.history_questions = main_test_record["questions"]
                        st.session_state.history_q_type = main_test_record["q_type"]
                        st.session_state.history_user_answers = main_test_record["user_answers"]
                        st.session_state.history_total_score = main_test_record["total_score"]
                        st.session_state.history_score_detail = main_test_record["score_detail"]
                with col_export:
                    if st.button("📤 导出测试结果", use_container_width=True):
                        level = "优秀" if main_test_record['total_score'] >=90 else "良好" if main_test_record['total_score'] >=80 else "中等" if main_test_record['total_score'] >=70 else "及格" if main_test_record['total_score'] >=60 else "不及格"
                        export_text = f"试卷：{current_exam['doc_title']}\n测试时间：{main_test_record['test_time']}\n总分：{main_test_record['total_score']}/100 分\n等级：{level}\n\n详细批改：\n"
                        for detail in main_test_record["score_detail"]:
                            export_text += f"Q{detail['index']}：总分{detail['total_score']}分，得分{detail['get_score']}分，{detail['remark']}\n"
                        st.download_button(
                            label="下载结果文件",
                            data=export_text,
                            file_name=f"历史试卷测试结果_{main_test_record['total_score']}分.txt",
                            mime="text/plain",
                            use_container_width=True
                        )
                st.divider()

                if st.session_state.history_user_answers:
                    st.subheader("📝 答题详情")
                    total_num = len(st.session_state.history_questions)
                    for i in range(total_num):
                        if i not in st.session_state.history_user_answers:
                            continue
                        data = st.session_state.history_user_answers[i]
                        title = data["title"]
                        user_ans = data["user_ans"]
                        correct_ans = data["correct_ans"]
                        analysis = data["analysis"]
                        opts = data.get("opts", [])
                        detail = next((d for d in st.session_state.history_score_detail if d["index"] == i+1), None)
                        
                        with st.container(border=True):
                            st.markdown(f"**Q{i+1}. {title}**")
                            
                            if st.session_state.history_q_type in ["单选", "多选"] and len(opts) == 4:
                                st.markdown(f"A. {opts[0]}")
                                st.markdown(f"B. {opts[1]}")
                                st.markdown(f"C. {opts[2]}")
                                st.markdown(f"D. {opts[3]}")
                                st.markdown("<br>", unsafe_allow_html=True)
                            
                            if detail:
                                st.markdown(f"**本题总分：{detail['total_score']}分 | 你的得分：{detail['get_score']}分**")
                            
                            col1, col2 = st.columns(2)
                            with col1:
                                user_display = format_answer(user_ans)
                                st.markdown(f"<p style='color: #ff4444; font-weight: bold;'>你的答案：{user_display}</p>", unsafe_allow_html=True)
                            with col2:
                                st.markdown(f"<p style='color: #0066cc; font-weight: bold;'>正确答案：{correct_ans}</p>", unsafe_allow_html=True)
                            
                            if detail:
                                status = detail["status"]
                                if status == "full":
                                    st.success(f"✅ 回答正确！{detail['remark']}")
                                elif status == "partial":
                                    st.warning(f"⚠️ 部分正确 | {detail['remark']}")
                                else:
                                    st.error(f"❌ 回答错误 | {detail['remark']}")
                            
                            with st.expander("📌 查看详细解析", expanded=False):
                                formatted_analysis = analysis.replace('<br>', '\n').replace('\n', '<br>')
                                st.markdown(f"<p style='color: #008000;'><strong>详细解析：</strong><br>{formatted_analysis}</p>", unsafe_allow_html=True)
                        
                        st.markdown("<br>", unsafe_allow_html=True)
                st.divider()
            
            st.subheader("📌 题目与答案")
            render_view_mode(current_exam["questions"], current_exam["q_type"])
        else:
            st.error("❌ 未找到该试卷，请返回历史试卷列表")
            if st.button("← 返回历史试卷列表"):
                st.session_state.history_current_exam_id = ""
                st.rerun()
    else:
        if not st.session_state.exam_history:
            st.info("暂无试题历史记录，请先生成试题")
        else:
            for idx, item in enumerate(st.session_state.exam_history):
                exam_id = item["id"]
                q_type = item["q_type"]
                doc_title = item["doc_title"]
                create_time = item["time"]
                main_test_record = item["main_test_record"]
                is_tested = main_test_record is not None
                test_tag = "✅ 已测试" if is_tested else "⚠️ 未测试"
                
                panel_title = f"{test_tag} | 试卷 {exam_id} | {q_type} | {doc_title} | {create_time}"
                with st.expander(panel_title, expanded=False):
                    if st.button(f"❌ 删除本组", key=f"del_exam_{exam_id}"):
                        st.session_state.exam_history.pop(idx)
                        st.success("✅ 本组试题已删除")
                        st.rerun()
                    
                    st.divider()
                    
                    if st.button("📋 查看试卷详情", use_container_width=True, key=f"view_exam_{exam_id}"):
                        st.session_state.history_current_exam_id = exam_id
                        st.session_state.history_questions = item["questions"]
                        st.session_state.history_q_type = item["q_type"]
                        st.session_state.history_user_answers = {}
                        st.session_state.history_total_score = 0
                        st.session_state.history_score_detail = []
                        st.rerun()

# 4. 知识点分析页面
elif st.session_state.page_state == "knowledge_analysis":
    render_knowledge_analysis_page()

# 5. 我的错题本页面
elif st.session_state.page_state == "wrong_book":
    st.title("❌ 我的错题本")
    st.divider()

    wrong_questions = st.session_state.wrong_questions
    unique_wrongs = []
    seen_titles = set()
    for q in wrong_questions:
        if q["question_title"] not in seen_titles:
            seen_titles.add(q["question_title"])
            unique_wrongs.append(q)
    wrong_questions = unique_wrongs

    if not wrong_questions:
        st.info("🎉 太棒了！你目前没有错题，继续保持~")
    else:
        total_wrong = len(wrong_questions)
        knowledge_count = len(set([kp for q in wrong_questions for kp in q["knowledge_points"]]))
        col1, col2 = st.columns(2)
        with col1:
            st.metric(label="累计错题数", value=total_wrong)
        with col2:
            st.metric(label="涉及知识点数", value=knowledge_count)
        st.divider()

        st.subheader("📝 错题重测")
        st.caption("勾选需要重新测试的错题，点击「开始重测」即可作答")
        selected_ids = []
        with st.container(border=True):
            for idx, q in enumerate(wrong_questions):
                q_id = f"wrong_{idx}"
                kp_text = "、".join(q["knowledge_points"]) if q["knowledge_points"] else "未分类知识点"
                checked = st.checkbox(
                    f"Q{idx+1} | {q['question_title']}（知识点：{kp_text}）",
                    key=q_id,
                    value=True
                )
                if checked:
                    selected_ids.append(idx)
        
        col_start, col_clear = st.columns([1,1])
        with col_start:
            if st.button("🚀 开始重测", use_container_width=True, type="primary"):
                if not selected_ids:
                    st.warning("请至少选择一道错题")
                else:
                    selected_questions = [wrong_questions[i] for i in selected_ids]
                    st.session_state.wrong_questions_for_test = selected_questions
                    st.session_state.wrong_exam_mode = "test"
                    st.session_state.wrong_user_answers = {}
                    st.session_state.wrong_graded = False
                    st.session_state.wrong_total_score = 0
                    st.session_state.wrong_score_detail = []
                    st.rerun()
        with col_clear:
            if st.button("🗑️ 清空全部错题", use_container_width=True):
                st.session_state.user_answer_records = [
                    record for record in st.session_state.user_answer_records
                    if record["is_correct"]
                ]
                st.session_state.wrong_questions = []
                st.success("✅ 全部错题已清空")
                st.rerun()
        st.divider()

        if st.session_state.wrong_exam_mode == "test":
            st.subheader("📝 错题重测作答")
            test_questions = st.session_state.wrong_questions_for_test
            total_num = len(test_questions)
            q_type = test_questions[0]["question_type"] if test_questions else "单选"

            submitted, updated_answers = render_test_mode(
                questions=[q["question_full_line"] for q in test_questions],
                current_q_type=q_type,
                total_question_num=total_num,
                key_prefix="wrong",
                user_answers=st.session_state.wrong_user_answers
            )
            
            st.session_state.wrong_user_answers = updated_answers

            if submitted:
                total_score, score_detail = grade_exam_answers(
                    user_answers=st.session_state.wrong_user_answers,
                    current_q_type=q_type,
                    total_question_num=total_num,
                    model=get_ali_embedding
                )
                st.session_state.wrong_total_score = total_score
                st.session_state.wrong_score_detail = score_detail
                st.session_state.wrong_graded = True
                st.session_state.wrong_exam_mode = "graded"
                st.rerun()

        elif st.session_state.wrong_exam_mode == "graded":
            st.markdown(f"<h2 style='text-align: center;'>📊 重测结果：{st.session_state.wrong_total_score} / 100 分</h2>", unsafe_allow_html=True)
            score = st.session_state.wrong_total_score
            if score >= 90:
                level = "优秀"
                level_color = "#008000"
            elif score >= 80:
                level = "良好"
                level_color = "#32cd32"
            elif score >= 70:
                level = "中等"
                level_color = "#ffd700"
            elif score >= 60:
                level = "及格"
                level_color = "#ffa500"
            else:
                level = "不及格"
                level_color = "#ff4444"
            st.markdown(f"<h3 style='text-align: center; color: {level_color};'>等级：{level}</h3>", unsafe_allow_html=True)
            st.divider()

            st.subheader("📝 详细批改")
            test_questions = st.session_state.wrong_questions_for_test
            total_num = len(test_questions)
            q_type = test_questions[0]["question_type"] if test_questions else "单选"

            for i in range(total_num):
                if i not in st.session_state.wrong_user_answers:
                    continue
                data = st.session_state.wrong_user_answers[i]
                title = data["title"]
                user_ans = data["user_ans"]
                correct_ans = data["correct_ans"]
                analysis = data["analysis"]
                opts = data.get("opts", [])
                detail = next((d for d in st.session_state.wrong_score_detail if d["index"] == i+1), None)
                
                with st.container(border=True):
                    st.markdown(f"**Q{i+1}. {title}**")
                    
                    if q_type in ["单选", "多选"] and len(opts) == 4:
                        st.markdown(f"A. {opts[0]}")
                        st.markdown(f"B. {opts[1]}")
                        st.markdown(f"C. {opts[2]}")
                        st.markdown(f"D. {opts[3]}")
                        st.markdown("<br>", unsafe_allow_html=True)
                    
                    if detail:
                        st.markdown(f"**本题总分：{detail['total_score']}分 | 你的得分：{detail['get_score']}分**")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        user_display = format_answer(user_ans)
                        st.markdown(f"<p style='color: #ff4444; font-weight: bold;'>你的答案：{user_display}</p>", unsafe_allow_html=True)
                    with col2:
                        st.markdown(f"<p style='color: #0066cc; font-weight: bold;'>正确答案：{correct_ans}</p>", unsafe_allow_html=True)
                    
                    if detail:
                        status = detail["status"]
                        if status == "full":
                            st.success(f"✅ 回答正确！{detail['remark']}")
                        elif status == "partial":
                            st.warning(f"⚠️ 部分正确 | {detail['remark']}")
                        else:
                            st.error(f"❌ 回答错误 | {detail['remark']}")
                    
                    with st.expander("📌 查看详细解析", expanded=False):
                        formatted_analysis = analysis.replace('<br>', '\n').replace('\n', '<br>')
                        st.markdown(f"<p style='color: #008000;'><strong>详细解析：</strong><br>{formatted_analysis}</p>", unsafe_allow_html=True)
                
                st.markdown("<br>", unsafe_allow_html=True)
            
            st.divider()
            col_retest, col_export, col_back = st.columns([1,1,1])
            with col_retest:
                if st.button("🔄 重新测试", use_container_width=True, type="primary"):
                    st.session_state.wrong_exam_mode = "test"
                    st.session_state.wrong_user_answers = {}
                    st.session_state.wrong_graded = False
                    st.session_state.wrong_total_score = 0
                    st.session_state.wrong_score_detail = []
                    st.rerun()
            with col_export:
                if st.button("📤 导出重测结果", use_container_width=True):
                    export_text = f"错题重测结果\n总分：{st.session_state.wrong_total_score}/100 分\n等级：{level}\n\n详细批改：\n"
                    for detail in st.session_state.wrong_score_detail:
                        export_text += f"Q{detail['index']}：总分{detail['total_score']}分，得分{detail['get_score']}分，{detail['remark']}\n"
                    st.download_button(
                        label="下载结果文件",
                        data=export_text,
                        file_name=f"错题重测结果_{st.session_state.wrong_total_score}分.txt",
                        mime="text/plain",
                        use_container_width=True
                    )
            with col_back:
                if st.button("← 返回错题列表", use_container_width=True):
                    st.session_state.wrong_exam_mode = "view"
                    st.session_state.wrong_user_answers = {}
                    st.session_state.wrong_graded = False
                    st.rerun()

        else:
            st.subheader("📋 错题详情列表")
            for idx, q in enumerate(wrong_questions):
                kp_text = "、".join(q["knowledge_points"]) if q["knowledge_points"] else "未分类知识点"
                with st.expander(f"Q{idx+1} | 知识点：{kp_text}", expanded=False):
                    st.markdown(f"**题干：{q['question_title']}**")
                    st.markdown(f"**题型：{q['question_type']}**")
                    st.markdown(f"**你的得分：{q['user_score']}/{q['total_score']}分**")
                    st.markdown(f"**是否答对：❌ 答错**")
                    st.divider()
                    st.markdown("**💡 知识点巩固建议**")
                    st.info(f"本题核心考察「{kp_text}」，建议重新回顾对应知识点的课程内容，结合解析理解考点，再通过重测巩固掌握。")

# 全局底部版权信息
st.divider()
st.caption("💡 NLP智能题库系统 | 基于阿里云向量API语义匹配 | 在线作答 | 精准评分 | 自动批改")